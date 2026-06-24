#!/usr/bin/env python3
"""
Generate PPTs from word.txt - v3 with improved formatting:
- Section headers: colored, 14pt
- Body text: black, 11pt
- Cleaned extra line breaks and indentation
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

CN_NUMS = {
    '一': '01', '二': '02', '三': '03', '四': '04', '五': '05',
    '六': '06', '七': '07', '八': '08', '九': '09', '十': '10',
    '十一': '11', '十二': '12', '十三': '13', '十四': '14', '十五': '15',
    '十六': '16', '十七': '17', '十八': '18', '十九': '19', '二十': '20',
}

# Colors
COLOR_DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_POLICY = RGBColor(0x1B, 0x5C, 0x9E)     # blue for 政策依据 header
COLOR_RISK = RGBColor(0xCC, 0x33, 0x33)        # red for 预计风险 header
COLOR_SOLUTION = RGBColor(0x00, 0x80, 0x40)    # green for 解决方法 header
COLOR_DESC = RGBColor(0x33, 0x33, 0x33)        # dark gray for 风险描述 header


def read_file(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()


def clean_text(text):
    """Global text cleanup: remove 佳木斯 references and page numbers."""
    text = text.replace('佳木斯市税务局', '')
    text = text.replace('佳木斯市', '')
    text = text.replace('\f', '')
    # Remove standalone page numbers (alone on a line)
    text = re.sub(r'\n\s*\d{1,3}\s*\n', '\n', text)
    # Remove inline page numbers between CJK punctuation and next CJK char
    text = re.sub(r'([。；）\)])\s*\d{1,3}\s*([一-鿿【第（\(])', r'\1\2', text)
    # Normalize multiple blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def split_into_chapters(text):
    """Split word.txt into chapters using chapter header positions."""
    chapter_pattern = re.compile(
        r'第([一二三四五六七八九十]+)章\s*([^\s]+?风险指引)'
    )
    all_matches = list(chapter_pattern.finditer(text))

    content_matches = []
    for m in all_matches:
        after = text[m.end():m.end()+30]
        if '...' not in after:
            content_matches.append(m)

    chapters = []
    for i, m in enumerate(content_matches):
        cn_num = m.group(1)
        title = m.group(2).strip()
        start = m.start()
        end = content_matches[i+1].start() if i+1 < len(content_matches) else len(text)
        chapter_num = CN_NUMS.get(cn_num, cn_num)
        chapters.append((title, text[start:end], chapter_num, cn_num))
    return chapters


def extract_risk_points(chapter_content):
    """Extract individual risk points from chapter content."""
    risk_pattern = re.compile(r'【风险点\s*(\d+)】')
    matches = list(risk_pattern.finditer(chapter_content))
    if not matches:
        return []
    risk_points = []
    for i, m in enumerate(matches):
        start = m.start()
        end = matches[i+1].start() if i+1 < len(matches) else len(chapter_content)
        risk_points.append(chapter_content[start:end].strip())
    return risk_points


def parse_risk_point_sections(content):
    """Parse a risk point into its constituent sections."""
    sections = {}

    rp_num_match = re.match(r'【风险点\s*(\d+)】', content)
    sections['rp_num'] = rp_num_match.group(1) if rp_num_match else '?'

    # Title: from 【风险点 N】 to 【风险描述】
    title_match = re.match(r'【风险点\s*\d+】\s*(.*?)(?=【风险描述】)', content, re.DOTALL)
    if title_match:
        sections['title'] = re.sub(r'\s+', '', title_match.group(1).strip())
    else:
        sections['title'] = ''

    # Extract each section
    desc_match = re.search(r'【风险描述】\s*(.*?)(?=【政策依据】)', content, re.DOTALL)
    if desc_match:
        sections['risk_desc'] = clean_section_body(desc_match.group(1))

    policy_match = re.search(r'【政策依据】\s*(.*?)(?=【预计风险】)', content, re.DOTALL)
    if policy_match:
        sections['policy'] = clean_section_body(policy_match.group(1))

    expected_match = re.search(r'【预计风险】\s*(.*?)(?=【解决方法】)', content, re.DOTALL)
    if expected_match:
        sections['risk_expected'] = clean_section_body(expected_match.group(1))

    solution_match = re.search(r'【解决方法】\s*(.*?)$', content, re.DOTALL)
    if solution_match:
        sections['solution'] = clean_section_body(solution_match.group(1))

    return sections


def clean_section_body(text):
    """
    Clean section body text:
    - Remove page number artifacts
    - Normalize line breaks (join PDF-broken lines, keep real paragraph breaks)
    - No leading/trailing whitespace
    """
    text = text.strip()

    # Remove inline page numbers: 。19根据 -> 。根据
    text = re.sub(r'([。；）\)])\s*\d{1,3}\s*([一-鿿【第（\(])', r'\1\2', text)
    text = re.sub(r'^\d{1,3}\s*([一-鿿【])', r'\1', text)

    # Join lines that are broken by PDF layout (single newline within a paragraph)
    # A line that doesn't end with CJK punctuation and the next line starts with CJK char
    # means it was broken mid-sentence
    lines = text.split('\n')
    merged = []
    for line in lines:
        line = line.strip()
        if not line:
            merged.append('')  # keep paragraph breaks
            continue
        if merged and merged[-1] and not merged[-1].endswith(('。', '；', '：', '）', ')', '】')):
            # This line continues the previous one
            merged[-1] += line
        else:
            merged.append(line)

    # Remove empty lines at boundaries
    while merged and not merged[0]:
        merged.pop(0)
    while merged and not merged[-1]:
        merged.pop()

    # Rebuild text with natural paragraph breaks
    result = []
    for item in merged:
        if not item:
            result.append('')
        else:
            result.append(item)

    return '\n'.join(result)


def add_para(tf, text, font_size=Pt(11), bold=False, color=COLOR_BLACK,
             space_before=Pt(2), space_after=Pt(2), alignment=PP_ALIGN.LEFT):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.line_spacing = Pt(font_size.pt * 1.3)
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return p


def add_body_text(tf, text, font_size=Pt(11)):
    """Add body text (black), splitting long paragraphs naturally."""
    paragraphs = text.split('\n')
    for para in paragraphs:
        para = para.strip()
        if not para:
            # Empty line = visual spacing
            add_para(tf, '', font_size=Pt(4), color=COLOR_BLACK,
                     space_before=Pt(0), space_after=Pt(0))
            continue

        # If paragraph is very long, split on sentence boundaries
        if len(para) > 500:
            chunks = re.split(r'(?<=[。；）\)])', para)
            current = ''
            for chunk in chunks:
                if not chunk.strip():
                    continue
                if len(current) + len(chunk) < 500:
                    current += chunk
                else:
                    if current.strip():
                        add_para(tf, current.strip(), font_size=font_size, color=COLOR_BLACK)
                    current = chunk
            if current.strip():
                add_para(tf, current.strip(), font_size=font_size, color=COLOR_BLACK)
        else:
            add_para(tf, para, font_size=font_size, color=COLOR_BLACK)


def add_section_header(tf, text, color):
    """Add a colored section header."""
    add_para(tf, text, font_size=Pt(14), bold=True, color=color,
             space_before=Pt(8), space_after=Pt(3))


def create_chapter_ppt(chapter_title, risk_points, chapter_num, output_dir):
    """Create a PowerPoint presentation for a chapter."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ===== Title Slide =====
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_BLUE

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = chapter_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.733), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"共 {len(risk_points)} 个风险点"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.733), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "税费合规指引 (2.0版)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
    p3.alignment = PP_ALIGN.CENTER

    # ===== Content Slides =====
    for rp_idx, rp_content in enumerate(risk_points):
        slide = prs.slides.add_slide(blank_layout)
        rp_content = clean_text(rp_content)
        sections = parse_risk_point_sections(rp_content)

        # ---- Header bar ----
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.65))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DARK_BLUE
        header.line.fill.background()

        tf_h = header.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = Inches(0.4)
        tf_h.paragraphs[0].alignment = PP_ALIGN.LEFT
        run_h = tf_h.paragraphs[0].add_run()
        run_h.text = f"{chapter_title}  |  风险点 {sections.get('rp_num', '?')}"
        run_h.font.size = Pt(15)
        run_h.font.bold = True
        run_h.font.color.rgb = COLOR_WHITE

        # ---- Content area ----
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(6.4))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Risk point title (first paragraph)
        rp_title = sections.get('title', '')
        if rp_title:
            add_para(tf, rp_title, font_size=Pt(14), bold=True, color=COLOR_DARK_BLUE,
                     space_before=Pt(0), space_after=Pt(6))

        # Risk description
        desc = sections.get('risk_desc', '')
        if desc:
            add_section_header(tf, '【风险描述】', COLOR_DESC)
            add_body_text(tf, desc, Pt(11))

        # Policy basis
        policy = sections.get('policy', '')
        if policy:
            add_section_header(tf, '【政策依据】', COLOR_POLICY)
            add_body_text(tf, policy, Pt(11))

        # Expected risk
        expected = sections.get('risk_expected', '')
        if expected:
            add_section_header(tf, '【预计风险】', COLOR_RISK)
            add_body_text(tf, expected, Pt(11))

        # Solution
        solution = sections.get('solution', '')
        if solution:
            add_section_header(tf, '【解决方法】', COLOR_SOLUTION)
            add_body_text(tf, solution, Pt(11))

    # Save
    safe_title = chapter_title.replace('/', '-').replace('、', '-')
    filename = f"{chapter_num}_{safe_title}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    print(f"  Saved: {filename} ({len(risk_points)} slides)")
    return filepath


def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    word_file = os.path.join(base_dir, 'word.txt')
    output_dir = os.path.join(base_dir, 'ppts_output_v3')
    os.makedirs(output_dir, exist_ok=True)

    print("Reading word.txt...")
    text = read_file(word_file)

    print("Splitting into chapters...")
    chapters = split_into_chapters(text)
    print(f"Found {len(chapters)} chapters\n")

    total_rp = 0
    for title, content, chapter_num, cn_num in chapters:
        print(f"Chapter {chapter_num}: {title}")
        risk_points = extract_risk_points(content)
        print(f"  {len(risk_points)} risk points")
        total_rp += len(risk_points)
        if risk_points:
            create_chapter_ppt(title, risk_points, chapter_num, output_dir)
        else:
            print(f"  WARNING: No risk points found!")

    print(f"\n{'='*60}")
    print(f"All PPTs saved to: {output_dir}")
    print(f"Total: {total_rp} risk points across {len(chapters)} chapters")


if __name__ == '__main__':
    main()
