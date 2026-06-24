#!/usr/bin/env python3
"""
Generate PPTs from word.txt (manual copy-paste from PDF).
word.txt has more accurate content ordering than pdftotext output.
Each slide = one risk point. Removes "佳木斯市税务局" references.
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

CN_NUMS = {
    '一': '01', '二': '02', '三': '03', '四': '04', '五': '05',
    '六': '06', '七': '07', '八': '08', '九': '09', '十': '10',
    '十一': '11', '十二': '12', '十三': '13', '十四': '14', '十五': '15',
    '十六': '16', '十七': '17', '十八': '18', '十九': '19', '二十': '20',
}

def read_file(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()

def clean_text(text):
    """Clean up text: remove 佳木斯 references and page number artifacts."""
    text = text.replace('佳木斯市税务局', '')
    text = text.replace('佳木斯市', '')
    # Remove form feeds
    text = text.replace('\f', '')
    # Remove isolated page numbers (1-3 digits alone on a line or between CJK punctuation)
    text = re.sub(r'\n\s*\d{1,3}\s*\n', '\n', text)
    # Remove page numbers that appear between CJK characters (common in word.txt)
    # Pattern: CJK_char + whitespace + digits + whitespace + CJK_char
    text = re.sub(r'([。；）\)])\s*\d{1,3}\s*([【第（\(])', r'\1\2', text)
    text = re.sub(r'([。；）\)])\s*\d{1,3}\s*([^\d])', r'\1\2', text)
    # Normalize whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()

def split_into_chapters(text):
    """Split word.txt into chapters using actual chapter header positions."""
    # Find ALL chapter header positions in the content
    # Pattern matches both TOC and content headers
    chapter_pattern = re.compile(
        r'第([一二三四五六七八九十]+)章\s*([^\s]+?风险指引)'
    )

    all_matches = list(chapter_pattern.finditer(text))

    # Separate TOC entries (followed by dots) from content entries
    content_matches = []
    for m in all_matches:
        after = text[m.end():m.end()+30]
        if '...' not in after:
            content_matches.append(m)

    chapters = []
    for i, m in enumerate(content_matches):
        cn_num = m.group(1)
        title = m.group(2).strip()
        start = m.start()
        end = content_matches[i+1].start() if i+1 < len(content_matches) else len(text)
        chapter_num = CN_NUMS.get(cn_num, cn_num)
        chapters.append((title, text[start:end], chapter_num, cn_num))

    return chapters

def extract_risk_points(chapter_content):
    """Extract individual risk points from chapter content."""
    risk_pattern = re.compile(r'【风险点\s*(\d+)】')
    matches = list(risk_pattern.finditer(chapter_content))
    if not matches:
        return []
    risk_points = []
    for i, m in enumerate(matches):
        start = m.start()
        end = matches[i+1].start() if i+1 < len(matches) else len(chapter_content)
        risk_points.append(chapter_content[start:end].strip())
    return risk_points

def parse_risk_point_sections(content):
    """Parse a risk point into its constituent sections."""
    sections = {}

    # Extract risk point number
    rp_num_match = re.match(r'【风险点\s*(\d+)】', content)
    rp_num = rp_num_match.group(1) if rp_num_match else '?'
    sections['rp_num'] = rp_num

    # Extract the full title (from 【风险点 N】 to 【风险描述】)
    title_match = re.match(r'【风险点\s*\d+】\s*(.*?)(?=【风险描述】)', content, re.DOTALL)
    if title_match:
        full_title = re.sub(r'\s+', '', title_match.group(1).strip())
        sections['title'] = full_title
    else:
        sections['title'] = ''

    # 【风险描述】
    desc_match = re.search(r'【风险描述】\s*(.*?)(?=【政策依据】)', content, re.DOTALL)
    if desc_match:
        sections['risk_desc'] = clean_section_text(desc_match.group(1))

    # 【政策依据】
    policy_match = re.search(r'【政策依据】\s*(.*?)(?=【预计风险】)', content, re.DOTALL)
    if policy_match:
        sections['policy'] = clean_section_text(policy_match.group(1))

    # 【预计风险】
    expected_match = re.search(r'【预计风险】\s*(.*?)(?=【解决方法】)', content, re.DOTALL)
    if expected_match:
        sections['risk_expected'] = clean_section_text(expected_match.group(1))

    # 【解决方法】
    solution_match = re.search(r'【解决方法】\s*(.*?)$', content, re.DOTALL)
    if solution_match:
        sections['solution'] = clean_section_text(solution_match.group(1))

    return sections

def clean_section_text(text):
    """Clean up a section's text content - remove page numbers, normalize whitespace."""
    text = text.strip()
    # Remove page number artifacts: 1-3 digit numbers that appear between CJK punctuation
    # Pattern: CJK_punctuation + optional_space + 1-3digit + optional_space + CJK_char
    text = re.sub(r'([。；）\)])\s*\d{1,3}\s*([一-鿿【第（\(])', r'\1\2', text)
    # Also handle numbers at start of text followed by CJK
    text = re.sub(r'^\d{1,3}\s*([一-鿿【])', r'\1', text)
    # Normalize whitespace
    text = re.sub(r'\s+', '', text)
    # Add logical line breaks after periods (for readability in PPT)
    text = re.sub(r'([。；])([^\s])', r'\1\n\2', text)
    return text.strip()

def add_para(tf, text, font_size=Pt(10), bold=False, color=RGBColor(0x33, 0x33, 0x33),
             space_before=Pt(1), space_after=Pt(1)):
    """Add a paragraph to text frame."""
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return p

def add_long_text(tf, text, font_size, color):
    """Add long text, splitting into readable chunks."""
    # Break on Chinese periods and semicolons
    chunks = re.split(r'(?<=[。；）\)])', text)
    current = ''
    for chunk in chunks:
        if not chunk.strip():
            continue
        if len(current) + len(chunk) < 500:
            current += chunk
        else:
            if current.strip():
                add_para(tf, current.strip(), font_size=font_size, color=color)
            current = chunk
    if current.strip():
        add_para(tf, current.strip(), font_size=font_size, color=color)

def create_chapter_ppt(chapter_title, risk_points, chapter_num, output_dir):
    """Create a PowerPoint presentation for a chapter."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ===== Title Slide =====
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = chapter_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.733), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"共 {len(risk_points)} 个风险点"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.733), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "税费合规指引 (2.0版)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
    p3.alignment = PP_ALIGN.CENTER

    # ===== Content Slides =====
    for rp_idx, rp_content in enumerate(risk_points):
        slide = prs.slides.add_slide(blank_layout)
        rp_content = clean_text(rp_content)
        sections = parse_risk_point_sections(rp_content)

        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.7))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
        header.line.fill.background()

        tf_h = header.text_frame
        tf_h.word_wrap = True
        run_h = tf_h.paragraphs[0].add_run()
        run_h.text = f"{chapter_title}  |  风险点 {sections.get('rp_num', '?')}"
        run_h.font.size = Pt(16)
        run_h.font.bold = True
        run_h.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Content area
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12.133), Inches(6.3))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Risk point title
        rp_title = sections.get('title', '')
        if rp_title:
            add_para(tf, rp_title, Pt(13), True, RGBColor(0x1B, 0x3A, 0x5C), Pt(0), Pt(6))

        # Sections
        sections_data = [
            ('【风险描述】', sections.get('risk_desc', ''), RGBColor(0x33, 0x33, 0x33)),
            ('【政策依据】', sections.get('policy', ''), RGBColor(0x1B, 0x5C, 0x9E)),
            ('【预计风险】', sections.get('risk_expected', ''), RGBColor(0xCC, 0x33, 0x33)),
            ('【解决方法】', sections.get('solution', ''), RGBColor(0x00, 0x80, 0x40)),
        ]

        for sec_title, sec_content, sec_color in sections_data:
            if not sec_content.strip():
                continue
            add_para(tf, sec_title, Pt(12), True, sec_color, Pt(6), Pt(2))
            add_long_text(tf, sec_content, Pt(9.5), sec_color)

    # Save
    safe_title = chapter_title.replace('/', '-').replace('、', '-')
    filename = f"{chapter_num}_{safe_title}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    print(f"  Saved: {filename} ({len(risk_points)} slides)")
    return filepath

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    word_file = os.path.join(base_dir, 'word.txt')
    output_dir = os.path.join(base_dir, 'ppts_output_v2')
    os.makedirs(output_dir, exist_ok=True)

    print("Reading word.txt...")
    text = read_file(word_file)

    print("Splitting into chapters...")
    chapters = split_into_chapters(text)
    print(f"Found {len(chapters)} chapters\n")

    total_rp = 0
    for title, content, chapter_num, cn_num in chapters:
        print(f"Chapter {chapter_num}: {title}")
        risk_points = extract_risk_points(content)
        print(f"  {len(risk_points)} risk points")
        total_rp += len(risk_points)
        if risk_points:
            create_chapter_ppt(title, risk_points, chapter_num, output_dir)
        else:
            print(f"  WARNING: No risk points found!")

    print(f"\n{'='*60}")
    print(f"All PPTs saved to: {output_dir}")
    print(f"Total: {total_rp} risk points across {len(chapters)} chapters")

if __name__ == '__main__':
    main()
